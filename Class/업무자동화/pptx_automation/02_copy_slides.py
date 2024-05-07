import copy

from pptx import Presentation


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    
    # 슬라이드를 복사하는 메서드
    def copy_slide(self, from_slide_no=0, slide_layout_no=6): #layout6은 빈 슬라이드
        # from_slide_no: 복사할 슬라이드 번호
        from_slide = self.ppt_file.slides[from_slide_no]

        # slide_layout_no: 복사할 슬라이드의 레이아웃 번호
        # 일반적으로 6번이 빈 슬라이드
        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        # 복사할 슬라이드의 모든 요소를 복사
        for shape in from_slide.shapes:
            el = shape.element #copy1
            new_element = copy.deepcopy(el) #copy2
            to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst') #추가

    # 슬라이드를 여러 개 복사하는 메서드
    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        for _ in range(slide_cnt):
            self.copy_slide(from_slide_no=from_slide_no)

    # 파일을 저장하는 메서드
    def save(self, filename):
        self.ppt_file.save(filename)

    # 슬라이드의 shape 이름과 번호를 매핑하는 메서드
    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, shape in enumerate(slide.shapes):
            shape_map[shape.name] = i #번호로 저장해서 딕셔너리에 저장
        return shape_map
    
    # 특정 슬라이드의 shape 텍스트를 변경하는 메서드
    def change_text(self, slide_no, label_map): #딕셔너리로 lable_map넘겨받음
        slide = self.ppt_file.slides[slide_no] #변경하려는 슬라이드 가져옴
        shape_map = self.get_shape_map(slide_no)
        
        #shape_map : product_name -> shape index 반환(맵핑)

        for shape_name, text in label_map.items(): 
            shape_no = shape_map[shape_name]
            slide.shapes[shape_no].text = text


if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물조사표.pptx")

    # # 슬라이드 0을 복사하여 새 슬라이드를 만들고, 추가되었는지 확인합니다.
    # # 파일을 저장한 다음 열어서 확인하는 방법도 있지만 print_slide_shapes를 사용하면 더 빠릅니다.
    # ppt_al.print_slide_shapes(0)
    # ppt_al.copy_slide(0)
    # ppt_al.print_slide_shapes(1) #추가된 상태. 저장하지 않으면 실제 파일에 반영안됨

    # # 1000개를 만들고 싶은 경우 999개를 매개변수로 입력
    # ppt_al.duplicate_n_slides(999)
    # ppt_al.save("[Auto]재물조사표_copy_slide.pptx")

    # # 0번 슬라이드의 4번 shape의 텍스트를 변경
    ppt_al.print_slide_shapes(0)
    slide = ppt_al.ppt_file.slides[0]
    slide.shapes[3].text = "Hello, World!" #3번 인덱스에 있는 shpe 변경
    ppt_al.print_slide_shapes(0)

    # # 편집 - 선택 창에서 shape 이름을 확인 및 변경 가능
    ppt_al.copy_slide(0)
    ppt_al.change_text(0, {"product_name":"32인치 게이밍 모니터", "model_no":"MO11223344"})
    ppt_al.change_text(1, {"product_name":"40인치 거실용 TV", "model_no":"MO7890"})
    ppt_al.save("[Auto]재물조사표_change_text.pptx")
