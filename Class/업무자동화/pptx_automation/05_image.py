import copy
import os

import pandas as pd
from pptx import Presentation
import openpyxl
from pptx.util import Pt


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    def copy_slide(self, from_slide_no=0, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]

        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        img_dict = {}
        for shape in from_slide.shapes:
            # shape 이름에 Picture가 포함되어 있으면 이미지로 처리
            # ppt에서 그림파일은 기본적으로 Picture라는 이름 부여
            if "Picture" in shape.name:
                image_filename = shape.name + '.jpg'
                with open(image_filename, 'wb') as f:
                    f.write(shape.image.blob)
                img_dict[image_filename] = [shape.left, shape.top, shape.width, shape.height]
            else:
                el = shape.element
                new_element = copy.deepcopy(el)
                to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

        # 이미지 복사
        for key, value in img_dict.items():
            to_slide.shapes.add_picture(key, value[0], value[1], value[2], value[3])
            os.remove(key)

    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        for _ in range(slide_cnt):
            self.copy_slide(from_slide_no=from_slide_no)

    def save(self, filename):
        self.ppt_file.save(filename)

    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, shape in enumerate(slide.shapes):
            shape_map[shape.name] = i
        return shape_map

    def change_text(self, slide_no, label_map, font_size=30):
        slide = self.ppt_file.slides[slide_no]
        shape_map = self.get_shape_map(slide_no)

        for shape_name, text in label_map.items():
            shape_no = shape_map[shape_name]
            text_frame = slide.shapes[shape_no].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)



if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물조사표_로고.pptx")

    df = pd.read_excel("재물목록.xlsx")

    print("count:", df["product_name"].count())
    slide_cnt = df["product_name"].count()
    ppt_al.duplicate_n_slides(slide_cnt - 1)

    for i, row in df.iterrows():
        print(i, row['product_name'], row['model_no'])
        label_map = {
            "product_name":row["product_name"], "model_no":row["model_no"]
        }
        ppt_al.change_text(i, label_map, 18)
    ppt_al.save("[Auto]재물조사표_로고_403.pptx")
