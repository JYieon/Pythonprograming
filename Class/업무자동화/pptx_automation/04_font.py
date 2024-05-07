import copy
import pandas as pd
from pptx import Presentation
from pptx.util import Pt


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    def copy_slide(self, from_slide_no=0, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]

        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        for shape in from_slide.shapes:
            el = shape.element
            new_element = copy.deepcopy(el)
            to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        for _ in range(slide_cnt):
            self.copy_slide(from_slide_no=from_slide_no)

    def save(self, filename):
        self.ppt_file.save(filename)

    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, shape in enumerate(slide.shapes):
            shape_map[shape.name] = i
        return shape_map

    def change_text(self, slide_no, label_map, font_size=30):
        slide = self.ppt_file.slides[slide_no]
        shape_map = self.get_shape_map(slide_no)

        for shape_name, text in label_map.items():
            shape_no = shape_map[shape_name]
            # slide.shapes[shape_no].text = text
            # 텍스트만 변경하지 않고, 폰트 사이즈도 변경하는 코드
            text_frame = slide.shapes[shape_no].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)



if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물조사표.pptx")

    df = pd.read_excel("재물목록.xlsx")

    slide_cnt = df["product_name"].count()
    ppt_al.duplicate_n_slides(slide_cnt - 1)

    for i, row in df.iterrows():
        print(i, row['product_name'], row['model_no'])
        ppt_al.change_text(i, {"product_name":row["product_name"]}, 32)
        ppt_al.change_text(i, {"model_no":row["model_no"]}, 24)

    ppt_al.save("[Auto]재물조사표_403_2.pptx")

