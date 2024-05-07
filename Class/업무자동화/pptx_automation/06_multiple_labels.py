import copy
import os

import pandas as pd
from pptx import Presentation
import openpyxl
from pptx.util import Pt


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    def copy_slide(self, from_slide_no=0, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]

        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        img_dict = {}
        for shape in from_slide.shapes:
            if "Picture" in shape.name:
                image_filename = shape.name + '.jpg'
                with open(image_filename, 'wb') as f:
                    f.write(shape.image.blob)
                img_dict[image_filename] = [shape.left, shape.top, shape.width, shape.height]
            else:
                el = shape.element
                new_element = copy.deepcopy(el)
                to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

        for key, value in img_dict.items():
            to_slide.shapes.add_picture(key, value[0], value[1], value[2], value[3])
            os.remove(key)

    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        for _ in range(slide_cnt):
            self.copy_slide(from_slide_no=from_slide_no)

    def save(self, filename):
        self.ppt_file.save(filename)

    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, shape in enumerate(slide.shapes):
            shape_map[shape.name] = i
        return shape_map

    def change_text(self, slide_no, label_map, font_size=30):
        slide = self.ppt_file.slides[slide_no]
        shape_map = self.get_shape_map(slide_no)

        for shape_name, text in label_map.items():
            shape_no = shape_map[shape_name]
            text_frame = slide.shapes[shape_no].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            print(shape_name)
            run.text = text
            run.font.size = Pt(font_size)



if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("명찰표.pptx")

    df = pd.read_excel("명찰목록.xlsx")
    print("count:", df["person_name"].count())
    slide_cnt = df["person_name"].count()

    # 한페이지에 여러 라벨이 들어가므로 나눈 몫만큼 슬라이드 복사
    if slide_cnt % 6 == 0:
        ppt_al.duplicate_n_slides(slide_cnt // 6 - 1)
    else:
        ppt_al.duplicate_n_slides(slide_cnt // 6)

    for i, row in df.iterrows():
        # 0, 0, 0, 0, 1, 1, 1, 1, 2, 2, 2, 2, 3, 3, 3, 3, ...
        page_idx = i // 6
        # 1, 2, 3, 4, 1, 2, 3, 4, 1, 2, 3, 4, 1, 2, 3, 4, ...
        label_idx = (i % 6) + 1

        print(i, row['person_name'], row['number'])
        ppt_al.change_text(page_idx, {
            f"person_name_{label_idx}": row["person_name"],
            f"number_{label_idx}": str(row["number"]) #str을 따로 붙여줌... -> 번호에 문자가 없어서!!
        }, 23)

    ppt_al.save("[Auto]명찰표.pptx")

