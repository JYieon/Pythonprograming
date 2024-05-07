# 패키지 설치 필요: pip install python-pptx
from pptx import Presentation

class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        # 열리지 않는 경우, 파일이름이나 경로를 확인해보세요.
        # 또는 파일이 열려있다면 닫아주세요.
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no] #원하는 슬라이드를 slide0에 저장

        for shape in slide0.shapes: #슬라이드에 있는 객체들 순회(shapes)
            print(f"slide_no:{slide_no} {shape.text}")

if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물조사표.pptx")
    ppt_al.print_slide_shapes(0) #원하는 슬라이드 선택
    
    #빈 텍스트가 출력되는 건 사각형모양(text 없어서 빈 문자 출력)